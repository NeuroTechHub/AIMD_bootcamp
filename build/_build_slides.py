"""Generate Antonio's bootcamp talk deck.

Output: presentations/bootcamp_talk_AL.pptx
Re-run anytime: `python build/_build_slides.py`

Style mirrors `presentations/Neurotech Bootcamp prelude.pptx`:
dark header bar, white body, teal accent bar at the bottom.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ----- theme ----------------------------------------------------------------

INK_DARK = RGBColor(0x3A, 0x3A, 0x3A)     # header / titles in body
INK_BODY = RGBColor(0x44, 0x44, 0x44)     # body text
INK_MUTED = RGBColor(0x6A, 0x6A, 0x6A)    # subheads / captions
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x4A, 0x8D, 0x7E)       # teal accent bar
BULLET_GREY = RGBColor(0x55, 0x55, 0x55)

FONT_FAMILY = "Calibri"   # python-pptx ships no font; Calibri is everywhere
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(1.1)
ACCENT_H = Inches(0.18)

# Where extracted source images live (see build/_extract_pptx_assets.py).
_REPO = Path(__file__).resolve().parent.parent
SRC_IMG = _REPO / "presentations" / "sources" / "brain_chip_2024" / "images"
M4_ASSETS = _REPO / "modules" / "M4-phosphene-simulation" / "assets"
NTH_LOGO = _REPO / "build" / "assets" / "nth_logo.png"
LOGOS_DIR = _REPO / "build" / "assets" / "logos"
TEAM_DIR = _REPO / "build" / "assets" / "team"
QRS_DIR = _REPO / "build" / "assets" / "qrs"
SHOTS_DIR = _REPO / "build" / "assets" / "module_shots"

# Title-slide footer logos, left→right. White marks on the dark title background.
TITLE_FOOTER_LOGOS = [
    LOGOS_DIR / "adcorpus_white.png",
    LOGOS_DIR / "nth_wordmark_white.png",
    LOGOS_DIR / "tudelft_white.png",
]


# ----- low-level helpers ----------------------------------------------------


def _add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else None
    shape.shadow.inherit = False
    return shape


def _add_text(slide, x, y, w, h, text, *,
              size=18, bold=False, color=INK_BODY,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              font=FONT_FAMILY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _add_bullets(slide, x, y, w, h, items, *,
                 size=22, color=INK_BODY, line_spacing=1.35):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        text, level = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        bullet = "● " if level == 0 else "○ "
        run = p.add_run()
        run.text = bullet + text
        run.font.name = FONT_FAMILY
        run.font.size = Pt(size if level == 0 else size - 4)
        run.font.color.rgb = color if level == 0 else INK_MUTED
    return tb


def _add_image_centered(slide, path, *, y, max_w, max_h):
    """Place image centered horizontally in [0, SLIDE_W], scaled to fit max box."""
    from PIL import Image as PILImage  # lazy import; ships with python-pptx via pillow
    with PILImage.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    if max_w / max_h > ratio:
        h = max_h
        w = int(max_h * ratio)
    else:
        w = max_w
        h = int(max_w / ratio)
    x = (SLIDE_W - w) // 2
    return slide.shapes.add_picture(str(path), x, y, w, h)


def _add_qr_corner(slide, qr_path, label):
    """Stamp a QR + tiny italic label in the bottom-right corner. Sits high
    enough above the accent bar that it doesn't collide with the slide-number
    stamp (`i / N` at SLIDE_H - 0.55 in)."""
    if qr_path is None or not Path(qr_path).exists():
        return
    qr_size = Inches(0.9)
    x = SLIDE_W - qr_size - Inches(0.4)
    # Lift the QR so caption ends ~0.15in above the slide-number stamp.
    y = SLIDE_H - ACCENT_H - qr_size - Inches(0.8)
    slide.shapes.add_picture(str(qr_path), x, y, qr_size, qr_size)
    tb = _add_text(slide, x - Inches(0.4), y + qr_size, qr_size + Inches(0.8),
                   Inches(0.22), label, size=9, color=INK_MUTED,
                   align=PP_ALIGN.CENTER)
    tb.text_frame.paragraphs[0].runs[0].font.italic = True


def _add_credit(slide, text):
    _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.55),
              SLIDE_W - Inches(1.2), Inches(0.3),
              text, size=11, color=INK_MUTED)


def _stamp_slide_numbers(prs):
    """Add a small ``i / N`` page indicator to every slide except the title.

    Sits just above the accent bar in the bottom-right corner. Colour
    auto-adapts to dark vs light slide backgrounds (dark = full-slide rect
    added first by ``_chrome(dark_full=True)``)."""
    n = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        if i == 1:
            continue  # title slide stays unnumbered
        first = next(iter(slide.shapes), None)
        is_dark = first is not None and first.width == SLIDE_W and first.height == SLIDE_H
        color = RGBColor(0xCC, 0xCC, 0xCC) if is_dark else INK_MUTED
        tb = slide.shapes.add_textbox(
            SLIDE_W - Inches(1.1), SLIDE_H - Inches(0.55),
            Inches(0.85), Inches(0.3),
        )
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{i} / {n}"
        run.font.name = FONT_FAMILY
        run.font.size = Pt(10)
        run.font.color.rgb = color


def _add_nth_logo(slide):
    """Place the NTH brand mark in the top-right corner, inside the dark header
    band. Sits at the same height regardless of dark_full vs. header-only chrome
    so the chrome reads consistently across body and divider slides."""
    if not NTH_LOGO.exists():
        return
    h = Inches(0.7)  # logo height; fits within the 1.1in header band
    from PIL import Image as PILImage
    with PILImage.open(NTH_LOGO) as im:
        iw, ih = im.size
    w = int(h * (iw / ih))
    x = SLIDE_W - w - Inches(0.35)
    y = (HEADER_H - h) // 2
    slide.shapes.add_picture(str(NTH_LOGO), x, y, w, h)


def _chrome(slide, title=None, *, dark_full=False, nth_logo=True):
    """Add the dark header bar, accent bar, and optional title.

    `nth_logo=True` (default) places the NTH brand mark top-right on every
    slide. Pass False for slides that should be unbranded (the talk's own
    title slide)."""
    if dark_full:
        _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, INK_DARK)
    else:
        _add_rect(slide, 0, 0, SLIDE_W, HEADER_H, INK_DARK)
        _add_rect(slide, 0, SLIDE_H - ACCENT_H, SLIDE_W, ACCENT_H, ACCENT)
    if title:
        # Reserve ~1.4in on the right so the title text doesn't collide with
        # the NTH logo when one is present.
        title_w = SLIDE_W - Inches(1.2) - (Inches(1.4) if nth_logo else Inches(0))
        _add_text(
            slide, Inches(0.6), Inches(0.25), title_w, Inches(0.7),
            title, size=30, bold=False, color=PAPER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    if nth_logo:
        _add_nth_logo(slide)


# ----- slide builders -------------------------------------------------------


def slide_title(prs, big, sub, byline=None, footer_logos=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Title slide is intentionally unbranded — the NTH chrome starts on slide 2.
    _chrome(s, dark_full=True, nth_logo=False)
    _add_rect(s, 0, SLIDE_H - ACCENT_H, SLIDE_W, ACCENT_H, ACCENT)
    _add_text(s, Inches(1), Inches(2.0), SLIDE_W - Inches(2), Inches(1.2),
              big, size=60, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(1), Inches(3.3), SLIDE_W - Inches(2), Inches(0.7),
              sub, size=28, color=PAPER, align=PP_ALIGN.CENTER)
    if byline:
        tb = _add_text(s, Inches(1), Inches(4.1), SLIDE_W - Inches(2), Inches(0.5),
                       byline, size=18, color=PAPER, align=PP_ALIGN.CENTER)
        # Italicise the byline run.
        tb.text_frame.paragraphs[0].runs[0].font.italic = True
    if footer_logos:
        _add_footer_logos(s, footer_logos)


def _add_footer_logos(slide, paths, *, row_h=Inches(0.75), gap=Inches(0.8)):
    """Lay out a horizontal row of logos centred along the bottom of the slide,
    sitting just above the accent bar. Each logo is scaled to the same height;
    its width follows the source aspect ratio."""
    from PIL import Image as PILImage
    widths = []
    for p in paths:
        if not Path(p).exists():
            widths.append(0)
            continue
        with PILImage.open(p) as im:
            iw, ih = im.size
        widths.append(int(row_h * (iw / ih)))
    total = sum(widths) + gap * max(0, len(paths) - 1)
    x = (SLIDE_W - total) // 2
    y = SLIDE_H - ACCENT_H - row_h - Inches(0.4)
    for p, w in zip(paths, widths):
        if w == 0:
            continue
        slide.shapes.add_picture(str(p), x, y, w, row_h)
        x += w + gap


def slide_section(prs, label, kicker=None, image=None, image_credit=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, dark_full=True)
    _add_rect(s, 0, SLIDE_H - ACCENT_H, SLIDE_W, ACCENT_H, ACCENT)
    if image is None:
        if kicker:
            _add_text(s, Inches(1), Inches(2.6), SLIDE_W - Inches(2), Inches(0.5),
                      kicker, size=20, color=ACCENT, align=PP_ALIGN.CENTER)
        _add_text(s, Inches(1), Inches(3.1), SLIDE_W - Inches(2), Inches(1.3),
                  label, size=54, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
        return
    # Hero-image layout: text on the left half, image on the right half.
    text_w = Inches(6.4)
    if kicker:
        _add_text(s, Inches(0.6), Inches(3.0), text_w, Inches(0.5),
                  kicker, size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.6), Inches(3.5), text_w, Inches(2.0),
              label, size=44, bold=True, color=PAPER, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.TOP)
    img_box_x = Inches(7.2)
    img_box_w = SLIDE_W - img_box_x - Inches(0.5)
    img_box_h = Inches(6.0)
    pic = _add_image_centered(s, SRC_IMG / image,
                              y=(SLIDE_H - img_box_h) // 2,
                              max_w=img_box_w, max_h=img_box_h)
    # Re-center horizontally inside the right column (override the slide-wide
    # centering that _add_image_centered does).
    pic.left = img_box_x + (img_box_w - pic.width) // 2
    if image_credit:
        _add_text(s, img_box_x, SLIDE_H - Inches(0.55), img_box_w, Inches(0.3),
                  image_credit, size=10, color=RGBColor(0xAA, 0xAA, 0xAA),
                  align=PP_ALIGN.CENTER)


def slide_bullets(prs, title, items, *, subhead=None, size=22):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title=title)
    y = Inches(1.7)
    if subhead:
        _add_text(s, Inches(0.8), y, SLIDE_W - Inches(1.6), Inches(0.5),
                  subhead, size=18, color=INK_MUTED)
        y = Inches(2.2)
    _add_bullets(s, Inches(0.8), y, SLIDE_W - Inches(1.6), Inches(5),
                 items, size=size)


def slide_module(prs, code, name, blurb, points):
    """Module spotlight: big M-code, name, blurb, supporting bullets."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title=f"{code} — {name}")
    _add_text(s, Inches(0.8), Inches(1.55), SLIDE_W - Inches(1.6), Inches(0.6),
              blurb, size=20, color=INK_MUTED)
    _add_bullets(s, Inches(0.8), Inches(2.4), SLIDE_W - Inches(1.6), Inches(4.6),
                 points, size=20)


def slide_image_focus(prs, title, image=None, caption=None, credit=None,
                      max_h=Inches(4.8), *, image_path=None):
    """Big image as the slide's hero; optional caption above and credit below.

    `image` is a filename inside SRC_IMG; pass `image_path=Path(...)` to use a
    file from anywhere else (e.g. M4_ASSETS for animated GIFs).
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title=title)
    y = Inches(1.55)
    if caption:
        _add_text(s, Inches(0.8), y, SLIDE_W - Inches(1.6), Inches(0.5),
                  caption, size=20, color=INK_MUTED)
        y = Inches(2.15)
    src = image_path if image_path is not None else SRC_IMG / image
    _add_image_centered(s, src, y=y,
                        max_w=SLIDE_W - Inches(1.6), max_h=max_h)
    if credit:
        _add_credit(s, credit)


def slide_video_focus(prs, title, video_path, *, poster_path=None,
                      caption=None, credit=None, max_h=Inches(4.4)):
    """Hero video. Plays in PowerPoint slideshow mode; static poster shows in
    PDF/print. Pick a representative mid-video frame for the poster."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title=title)
    y = Inches(1.55)
    if caption:
        _add_text(s, Inches(0.8), y, SLIDE_W - Inches(1.6), Inches(0.7),
                  caption, size=20, color=INK_MUTED)
        y = Inches(2.35)
    if poster_path is not None and Path(poster_path).exists():
        from PIL import Image as PILImage
        with PILImage.open(poster_path) as im:
            iw, ih = im.size
        ratio = iw / ih
    else:
        ratio = 16 / 9
    max_w = SLIDE_W - Inches(1.6)
    if max_w / max_h > ratio:
        h = max_h
        w = int(max_h * ratio)
    else:
        w = max_w
        h = int(max_w / ratio)
    x = (SLIDE_W - w) // 2
    s.shapes.add_movie(
        str(video_path), x, y, w, h,
        poster_frame_image=str(poster_path) if poster_path else None,
        mime_type="video/mp4",
    )
    if credit:
        _add_credit(s, credit)


def slide_bullets_image(prs, title, items, image=None, credit=None,
                        subhead=None, image_h=Inches(3.2),
                        image_path=None, qr_path=None, qr_label=None):
    """Compressed bullets at top, wide image below — for figure-heavy beats.

    Reserves vertical space for wrapped lines: long bullets get a 2-line
    budget so the image below doesn't ride up onto a wrapped second line.
    Heuristic: ~95 chars per line at the current 18 pt font and slide width.

    Pass ``image_path=Path(...)`` to load an image from outside SRC_IMG
    (e.g. a module HTML screenshot). Pass ``qr_path`` + ``qr_label`` to
    stamp a QR code with an italic caption in the bottom-right corner."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title=title)
    y = Inches(1.55)
    if subhead:
        _add_text(s, Inches(0.8), y, SLIDE_W - Inches(1.6), Inches(0.5),
                  subhead, size=18, color=INK_MUTED)
        y = Inches(2.05)
    chars_per_line = 95
    line_h = Inches(0.4)
    bullet_h = sum(max(1, -(-len(item) // chars_per_line)) for item in items) * line_h
    _add_bullets(s, Inches(0.8), y, SLIDE_W - Inches(1.6), bullet_h,
                 items, size=18, line_spacing=1.2)
    img_y = y + bullet_h + Inches(0.15)
    src = image_path if image_path is not None else SRC_IMG / image
    # Narrow the image's max width if a QR is present so the corner stays clear.
    max_w = SLIDE_W - Inches(2.6) if qr_path else SLIDE_W - Inches(1.6)
    _add_image_centered(s, src, y=img_y, max_w=max_w, max_h=image_h)
    if credit:
        _add_credit(s, credit)
    if qr_path:
        _add_qr_corner(s, qr_path, qr_label or "scan to open")


def slide_three_columns(prs, title, columns, *, subhead=None):
    """Three labeled columns of bullets. Each column = (label, [bullets])."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title=title)
    y = Inches(1.55)
    if subhead:
        _add_text(s, Inches(0.8), y, SLIDE_W - Inches(1.6), Inches(0.5),
                  subhead, size=18, color=INK_MUTED)
        y = Inches(2.15)
    margin = Inches(0.8)
    gap = Inches(0.35)
    n = len(columns)
    col_w = (SLIDE_W - margin * 2 - gap * (n - 1)) / n
    # Use all the space between subhead and the accent bar — the longest
    # column may carry six bullets.
    col_h = SLIDE_H - (y + Inches(0.55)) - Inches(0.35)
    for i, (label, bullets) in enumerate(columns):
        x = margin + (col_w + gap) * i
        _add_text(s, x, y, col_w, Inches(0.5),
                  label, size=22, bold=True, color=ACCENT)
        _add_bullets(s, x, y + Inches(0.55), col_w, col_h,
                     bullets, size=16, line_spacing=1.25)


def slide_logos(prs, title, columns, *, subhead=None):
    """3-column 'who's building' grid: light tiles holding program logos.

    Each column is ``(header, entries)`` where each entry is
    ``(short_label, caption, logo_path_or_None)``. ``logo_path=None`` renders
    the label centred inside the tile (text-only); otherwise the logo is
    placed inside the tile and the label appears as the caption."""
    from PIL import Image as PILImage
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title=title)
    y = Inches(1.55)
    if subhead:
        _add_text(s, Inches(0.8), y, SLIDE_W - Inches(1.6), Inches(0.5),
                  subhead, size=18, color=INK_MUTED)
        y = Inches(2.15)
    margin = Inches(0.8)
    gap = Inches(0.35)
    n = len(columns)
    col_w = (SLIDE_W - margin * 2 - gap * (n - 1)) / n
    tile_h = Inches(0.7)
    cap_h = Inches(0.28)
    row_gap = Inches(0.05)
    row_h = tile_h + cap_h + row_gap
    pad = Inches(0.1)
    inner_w = col_w - pad * 2
    inner_h = tile_h - pad * 2
    tile_fill = RGBColor(0xF2, 0xF2, 0xF2)
    for i, (header, entries) in enumerate(columns):
        x = margin + (col_w + gap) * i
        _add_text(s, x, y, col_w, Inches(0.4),
                  header, size=20, bold=True, color=ACCENT,
                  align=PP_ALIGN.CENTER)
        ey = y + Inches(0.55)
        for label, caption, logo in entries:
            _add_rect(s, x, ey, col_w, tile_h, tile_fill)
            if logo is not None and Path(logo).exists():
                with PILImage.open(logo) as im:
                    iw, ih = im.size
                ratio = iw / ih
                if inner_w / inner_h > ratio:
                    h = inner_h
                    w = int(inner_h * ratio)
                else:
                    w = inner_w
                    h = int(inner_w / ratio)
                lx = x + (col_w - w) // 2
                ly = ey + (tile_h - h) // 2
                s.shapes.add_picture(str(logo), lx, ly, w, h)
            else:
                _add_text(s, x, ey, col_w, tile_h, label,
                          size=22, bold=True, color=INK_DARK,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text(s, x, ey + tile_h + Inches(0.02),
                      col_w, cap_h, caption,
                      size=10, color=INK_MUTED, align=PP_ALIGN.CENTER)
            ey += row_h


def slide_bullets_side_image(prs, title, items, image=None, *,
                             subhead=None, image_credit=None,
                             image_path=None, size=18):
    """Bullets on the left ~60%, an image on the right ~35%."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title=title)
    y = Inches(1.55)
    text_w = Inches(7.4)
    if subhead:
        _add_text(s, Inches(0.6), y, text_w, Inches(0.5),
                  subhead, size=18, color=INK_MUTED)
        y = Inches(2.1)
    _add_bullets(s, Inches(0.6), y, text_w, Inches(5),
                 items, size=size)
    img_x = Inches(8.2)
    img_w = SLIDE_W - img_x - Inches(0.6)
    img_h = Inches(4.6)
    img_y = (SLIDE_H - img_h) // 2 + Inches(0.4)
    src = image_path if image_path is not None else SRC_IMG / image
    pic = _add_image_centered(s, src, y=img_y, max_w=img_w, max_h=img_h)
    pic.left = img_x + (img_w - pic.width) // 2
    if image_credit:
        _add_text(s, img_x, SLIDE_H - Inches(0.55), img_w, Inches(0.3),
                  image_credit, size=10, color=RGBColor(0xAA, 0xAA, 0xAA),
                  align=PP_ALIGN.CENTER)


def slide_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title="Today's pipeline")
    _add_text(s, Inches(0.8), Inches(1.55), SLIDE_W - Inches(1.6), Inches(0.6),
              "Five modules, end-to-end — image in, perception out, loop closed.",
              size=20, color=INK_MUTED)
    labels = [
        ("M1", "Computer\nvision"),
        ("M2", "Gaze &\nDeepGaze"),
        ("M3", "Neuromodulation\n& stim"),
        ("M4", "Phosphene\nsimulation"),
        ("M5", "Decoding &\nclosed loop"),
    ]
    n = len(labels)
    margin = Inches(0.8)
    gap = Inches(0.25)
    box_w = (SLIDE_W - margin * 2 - gap * (n - 1)) / n
    box_h = Inches(2.6)
    y = Inches(3.1)
    for i, (code, name) in enumerate(labels):
        x = margin + (box_w + gap) * i
        _add_rect(s, x, y, box_w, box_h, INK_DARK)
        _add_text(s, x, y + Inches(0.3), box_w, Inches(0.8),
                  code, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _add_text(s, x, y + Inches(1.2), box_w, Inches(1.2),
                  name, size=20, color=PAPER, align=PP_ALIGN.CENTER)
        if i < n - 1:
            # arrow chevron in the gap
            ax = x + box_w + Emu(20000)
            ay = y + box_h / 2 - Inches(0.12)
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, ay, gap - Inches(0.1), Inches(0.24)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()


def slide_team(prs, title, members, *, cols=None):
    """Team-portrait grid: circular-cropped headshots with a name caption.

    `members` is a list of ``(display_name, image_path)`` tuples. ``cols``
    controls the row width; defaults to 3 when len>=4, 2 otherwise. Rows
    of fewer portraits centre themselves."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _chrome(s, title=title)

    n = len(members)
    if cols is None:
        cols = 3 if n >= 4 else n
    rows = (n + cols - 1) // cols

    margin_x = Inches(0.8)
    avail_w = SLIDE_W - 2 * margin_x
    cell_w = avail_w // cols
    portrait = min(cell_w - Inches(0.4), Inches(2.0))  # cap so labels fit
    row_h = portrait + Inches(0.9)                     # portrait + name caption
    top_y = Inches(1.7)
    avail_h = SLIDE_H - top_y - Inches(0.7)
    if rows > 1 and row_h * rows > avail_h:
        row_h = avail_h // rows
        portrait = row_h - Inches(0.9)

    for i, (name, img) in enumerate(members):
        r = i // cols
        c = i % cols
        members_this_row = min(cols, n - r * cols)
        row_total_w = cell_w * members_this_row
        row_x_start = (SLIDE_W - row_total_w) // 2
        cx = row_x_start + cell_w * c + (cell_w - portrait) // 2
        cy = top_y + r * row_h
        _add_circular_portrait(s, img, cx, cy, portrait)
        _add_text(
            s, row_x_start + cell_w * c, cy + portrait + Inches(0.1),
            cell_w, Inches(0.7),
            name, size=16, color=INK_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )


def _add_circular_portrait(slide, image_path, x, y, size):
    """Place an image as a centred circular crop of side `size`.

    Centre-crops the source to a square (so faces aren't stretched), then
    swaps the picture's preset geometry from rect → ellipse to clip it
    inside a circle."""
    from PIL import Image as PILImage
    from pptx.oxml.ns import qn

    image_path = Path(image_path)
    if not image_path.exists():
        return None
    with PILImage.open(image_path) as im:
        iw, ih = im.size
    pic = slide.shapes.add_picture(str(image_path), x, y, size, size)
    if ih > iw:
        cut = (ih - iw) / ih / 2
        pic.crop_top = cut
        pic.crop_bottom = cut
    elif iw > ih:
        cut = (iw - ih) / iw / 2
        pic.crop_left = cut
        pic.crop_right = cut
    spPr = pic._element.find(qn("p:spPr"))
    if spPr is not None:
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            prstGeom.set("prst", "ellipse")
    return pic


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # No NTH logo on the sign-off — the call-to-action stays the only thing on screen.
    _chrome(s, dark_full=True, nth_logo=False)
    _add_rect(s, 0, SLIDE_H - ACCENT_H, SLIDE_W, ACCENT_H, ACCENT)
    _add_text(s, Inches(1), Inches(2.6), SLIDE_W - Inches(2), Inches(1.0),
              "Let's build.", size=60, bold=True, color=PAPER,
              align=PP_ALIGN.CENTER)
    _add_text(s, Inches(1), Inches(3.9), SLIDE_W - Inches(2), Inches(0.7),
              "Pizza, prizes, and the NTH GitHub waiting for your repo.",
              size=22, color=PAPER, align=PP_ALIGN.CENTER)
    # Pipeline-foundation credit (mirrors the HTML module footers).
    _add_text(s, Inches(1), SLIDE_H - Inches(0.95), SLIDE_W - Inches(2), Inches(0.4),
              "Pipeline foundation:  Lozano et al. 2020 — Neurolight  ·  "
              "Int. J. Neural Syst. 30(09):2050045  ·  doi:10.1142/S0129065720500458",
              size=11, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)


# ----- the deck -------------------------------------------------------------


def build(out_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # The prelude deck covers the NTH/organizing-team intro, vision, program,
    # hands-on framing, tracks, and sponsors. This deck opens with the talk's
    # own title slide and then dives in.

    # 0 — Title slide for this talk.
    slide_title(
        prs, "Cortical visual prostheses",
        "Brain Writing with the Neurotech Hub",
        byline="Antonio Lozano 2026",
        footer_logos=TITLE_FOOTER_LOGOS,
    )

    # 1 — Who's giving this talk (bullets on left, anatomy hero on right).
    slide_bullets_side_image(
        prs, "About me",
        subhead="Antonio Lozano. Neural engineer for vision.",
        size=16,
        items=[
            "Postdoc at UMH in Elche. I run the AI side of CORTIVIS — the first cortical visual prosthesis trial in human volunteers.",
            "Four years before that at the NIN in Amsterdam (Roelfsema lab). High-channel cortical implants, neural phosphene mapping, and the dynaphos simulator that powers this bootcamp's M4.",
            "Co-lead of INTENSE, a Dutch consortium pulling translational neurotech toward the clinic. My work package is blindness.",
            "Senior scientist at Ruten Inc. on the side. BCI R&D and the kind of CV pipelines you'll meet in M1.",
            "PhD in 2022 from UPCT in Cartagena (Cum Laude): \"AI-endowed visual neuroprosthesis for the blind.\"",
            "What I actually do day to day: sit between experimental neuroscience, deep learning, and a clinical trial that has to keep working tomorrow.",
        ],
        image="s01_p01_0291ff0f.png",
        image_credit="from A. Lozano, Brain & the Chip II (Elche 2024)",
    )

    # 2 — One slide replaces the four definition slides + the section divider.
    #     Three columns of where the field actually is in 2026.
    # TODO (RV review 2026-05-31): this slide is text-heavy across three
    # columns. Consider folding in a small graphic (e.g. an implant-class
    # photo strip or the addressable-market chart) — would require a
    # two-column-plus-image layout, not a code-level tweak.
    slide_three_columns(
        prs, "Neurotechnology in 2026",
        subhead="Where the field is today — what's deployed, what's emerging, and where this bootcamp lives.",
        columns=[
            ("Deployed clinical", [
                "Cochlear implants — over a million people hearing",
                "Deep brain stim for Parkinson, dystonia, OCD (200k+ implants)",
                "Spinal-cord stim, and responsive neurostimulation (RNS) for epilepsy",
                "Mature read-and-write, but at low channel count",
            ]),
            ("Emerging clinical", [
                "BrainGate motor BCIs — trials for over a decade",
                "Synchron's endovascular Stentrode",
                "Neuralink's first human implantations",
                "CEA-Leti / WIMAGINE; Cortec Brain Interchange",
                "Higher channels + closed loops (stim responding to brain state)",
                "Morgan Stanley 2024: ~$400B market, <3% penetration by 2045",
            ]),
            ("Scope of this bootcamp", [
                "Cortical visual prostheses — restore sight by writing to V1",
                "Read what V1 needs, write the percept",
                "Camera → computer vision → stimulation → percept",
                "This bootcamp: one working piece per module",
            ]),
        ],
    )

    # 2 — Intracortical Visual Prostheses (the 18:10 talk)
    slide_section(prs, "Intracortical Visual Prostheses",
                  kicker="Why we're building this pipeline",
                  image="s58_p01_9c76d02b.jpg",
                  image_credit="Bernadeta Gómez — UMH / Fernández cohort")
    slide_bullets_image(
        prs, "The problem",
        items=["Profound blindness: millions of people, many causes, no treatment for most",
               "Retinal and optic-nerve implants need everything downstream still working",
               "Cortical stim skips everything upstream — addressing V1 directly",
               "Electrode design (penetrating vs surface): optimal intervention strategy is still open",
               "Cortical is invasive, but reaches the broadest blind population"],
        image="s39_p01_ee164db1.png",
        subhead="The visual world is complex. Our bandwidth to represent it is limited.",
        credit="from A. Lozano, Brain & the Chip II (Elche 2024)",
        image_h=Inches(2.6),
    )
    slide_bullets_image(
        prs, "The visual pathway",
        items=["Eye → LGN → V1 (and beyond)",
               "Retinotopy: neighbouring V1 neurons see neighbouring parts of the visual field",
               "Cortical magnification: a tiny patch of fovea takes up a huge slice of V1",
               "Which is why V1 is the most-studied target for an intracortical implant"],
        image="s15_p01_ed547d2f.png",
        subhead="V1 is a map of the visual field. That map is what we'll be writing to.",
        credit="from A. Lozano, Brain & the Chip II (Elche 2024)",
        image_h=Inches(2.6),
    )
    # Field-map slide: puts cortical in context against retina + LGN.
    # See docs/vision-restoration-field.md for the long-form reference
    # with primary-source links. Bullets kept tight (single line each)
    # — slide_three_columns can't reflow at this width.
    slide_three_columns(
        prs, "Where to intervene",
        subhead="Each tissue target reaches a different patient population. This bootcamp builds the V1 piece.",
        columns=[
            ("Retina", [
                "Photoreceptor loss (AMD/RP). Needs intact downstream path.",
                "PRIMA (Pixium → Science Corp): subretinal photovoltaic, NEJM 2025",
                "Argus II (Second Sight → Cortigent): epi-retinal; 350+; ended 2019",
                "BVT (Australia): suprachoroidal 44-ch; 2025 RP safety data",
                "GenSight GS030: optogenetic + NIR goggles; PIONEER",
            ]),
            ("LGN (thalamus)", [
                "Retinal + optic-nerve causes. Needs intact V1.",
                "SIGHTED (EIC, Phosphoenix-coordinated): >1000-ch LGN",
                "Builds on NeuraViPeR (concluded Feb 2025)",
                "Preclinical; preparing first-in-human",
            ]),
            ("V1 cortex — this bootcamp", [
                "Broadest population — anything upstream of V1.",
                "Orion (Cortigent): surface 60-ch; drawing on cortex (Cell 2020)",
                "CORTIVIS (Fernández/UMH Elche): penetrating 96-ch; Sci Adv 2025",
                "Neuralink Blindsight: threads; FDA Breakthrough Sept 2024",
                "ReVision Occular (KU Leuven, BE): thin-film; FDA Breakthrough 2026",
            ]),
        ],
    )

    # Cortical landscape: bullet-only slide (no portrait) so the six
    # programs each get a full line. Was hitting image-overlap clipping
    # under slide_bullets_image's vertical budget.
    slide_bullets(
        prs, "Cortical landscape",
        subhead="Crowded V1 cortex: surface + 3 penetrating approaches; thin-film and threads are the trend.",
        items=[
            "CORTIVIS (Fernández / UMH Elche): active first-in-human; Sci Adv 2025 — bidirectional implant in 2 volunteers (NEJM 2021, Gómez)",
            "Orion (Cortigent, ex–Second Sight): EFS; Argus-derived surface electrodes need ~mA → Beauchamp draws on cortex (Cell 2020)",
            "Neuralink Blindsight: penetrating threads; FDA Breakthrough Sept 2024; first humans planned 2026 (IEEE Spectrum: expect low-res early on)",
            "ReVision Implant Occular (Ceyssens & Janssen, KU Leuven, BE): ultra-thin-film flexible; FDA Breakthrough 2026; FIH Q4 2026",
            "Phosphoenix (NL spin-off of NIN): coordinates SIGHTED (LGN) + the concluded NeuraViPeR (V1)",
            "In 2026 the bottleneck isn't the electrode — it's mapping, decoding, and control",
        ],
        size=16,
    )

    # Who's building — logo grid by tissue. Visual recap of the field-map
    # text on the previous two slides. Logos sourced from each program's
    # public press kit / homepage (see docs/vision-restoration-field.md
    # for URLs). White wordmarks were inverted to dark variants so they
    # read on the light tiles.
    _P = LOGOS_DIR / "programs"
    slide_logos(
        prs, "Who's building",
        subhead="Active programs in 2026 by target tissue. Logos via each program's homepage / press kit.",
        columns=[
            ("Retina", [
                ("PRIMA",        "PRIMA · Science Corp · subretinal",          _P / "science_prima.png"),
                ("Argus II",     "Argus II · Cortigent · epi-retinal",         _P / "cortigent_dark.png"),
                ("GS030",        "GS030 · GenSight · optogenetics + NIR",      _P / "gensight_dark.png"),
                ("Bionic Sight", "Bionic Sight · Nirenberg · optogenetics",    _P / "bionicsight_dark.png"),
            ]),
            ("LGN (thalamus)", [
                ("SIGHTED",      "SIGHTED · EIC consortium · >1000-ch",        _P / "sighted.jpg"),
                ("NeuraViPeR",   "NeuraViPeR · EU H2020 · concluded Feb 2025", _P / "neuraviper.png"),
            ]),
            ("V1 cortex — this bootcamp", [
                ("CORTIVIS",     "Fernández / UMH Elche · penetrating 96-ch", None),
                ("Orion",        "Orion · Cortigent · surface 60-ch",          _P / "cortigent_dark.png"),
                ("Blindsight",   "Blindsight · Neuralink · threads · FDA BD 2024",   _P / "neuralink.png"),
                ("ReVision",     "ReVision · KU Leuven · thin-film · FDA BD 2026",   _P / "revision_implant_dark.png"),
            ]),
        ],
    )

    # 3 — Bridge: borrow the 4-walls framing from "Brain & the Chip II".
    slide_bullets(
        prs, "Four walls between us and a real vision implant",
        subhead="Four open problems. This bootcamp gives you a working module against each one.",
        items=[
            "Surgical planning: where do you actually put the implant for the best coverage?",
            ("M4 phosphene sim + vimplant2 — design and compare layouts in the browser", 1),
            "Phosphene mapping: each electrode produces a different percept. Which one?",
            ("M3 stim explorer, with Granley & Beyeler's temporal patterns", 1),
            "Computer vision in real scenes: tiny bandwidth, useful percepts",
            ("M1 front-end CV + M2 gaze and saliency", 1),
            "Human-in-the-loop stim: close the loop with the user, not the model",
            ("M5 decoding + the end-to-end differentiable prosthesis", 1),
            "Today: one working module per wall.",
        ])

    # 4 — Real phosphene simulation: the eLife Video 2 example, embedded as
    #     a movie. Plays in slideshow mode (click). Three panels: real scene,
    #     edge map, phosphene field. Poster frame at t=6s shows all three.
    slide_video_focus(
        prs, "What patients actually see",
        video_path=M4_ASSETS / "elife-85812-video2.mp4",
        poster_path=M4_ASSETS / "elife-85812-video2-poster.jpg",
        caption="Real scene → edges → phosphenes. Sparse, punctate, dynamic — not pixels. (Click to play in slideshow mode.)",
        credit="van der Grinten et al. 2024, eLife — Video 2 (CC BY, doi:10.7554/eLife.85812)",
    )

    # 4 — Stim & safety (divider cut for time; the title slide below carries it)
    slide_bullets_image(
        prs, "Stim parameters that matter",
        items=["Amplitude (µA): how bright the phosphene looks",
               "Pulse width (µs): with amplitude, sets the total charge per phase",
               "Frequency (Hz): drives flicker fusion and adaptation",
               "Train shape: how many pulses, how long, with what duty cycle"],
        image="s43_p01_30e792f5.png",
        subhead="Perception probability shifts with charge, pulse width, frequency, train duration.",
        credit="Fernández et al. — phosphene perception psychophysics (Moran cohort)",
        image_h=Inches(2.6),
    )
    slide_bullets(prs, "Safety limits",
                  ["Shannon-k: log(D) + k·log(Q/A) ≤ k_max (Cogan 2016)",
                   "Charge density per phase: the threshold above which tissue starts to damage",
                   "Total charge across a whole session",
                   "M3 enforces all of these live, before the stim button does anything"])

    # 5 — Pipeline tour (divider cut; the loop figure carries the section)
    slide_image_focus(
        prs, "Vision implant — the loop in one figure",
        image="s09_p01_6614a827.png",
        caption="Camera → image processor → cortical implant → percept. Every module today is one piece of this loop.",
        credit="from A. Lozano, Brain & the Chip II (Elche 2024)",
        max_h=Inches(4.6),
    )
    slide_pipeline(prs)
    # Each M1–M5 spotlight is duplicated: original (figure from Antonio's
    # Brain & the Chip talk) followed by a v2 with a live screenshot of the
    # corresponding HTML playground + a QR pointing to the hosted page.
    # User to pick which version to keep and delete the other.

    # --- M1 ---
    slide_bullets_image(
        prs, "M1 — Computer vision",
        items=["OpenCV.js in the browser: Sobel, Canny, thresholding",
               "YOLO + COCO-SSD via TF.js — live object detection from your webcam",
               "Five processing modes you can flip between live",
               "Owners: Lefteris & Jorge"],
        image="s51_p01_6c988087.png",
        subhead="Pixels in, features out. The front end of any prosthesis.",
        credit="real-scene → edges → sparse phosphenes (from A. Lozano, Brain & the Chip II)",
        image_h=Inches(2.1),
    )
    slide_bullets_image(
        prs, "M1 — Computer vision",
        items=["OpenCV.js in the browser: Sobel, Canny, thresholding",
               "YOLO + COCO-SSD via TF.js — live object detection from your webcam",
               "Five processing modes you can flip between live",
               "Owners: Lefteris & Jorge"],
        image_path=SHOTS_DIR / "M1.png",
        subhead="Pixels in, features out. The front end of any prosthesis.",
        credit="screenshot from modules/M1-computer-vision.html",
        image_h=Inches(2.8),
        qr_path=QRS_DIR / "M1.png",
        qr_label="scan ▶ M1 live",
    )

    # --- M2 ---
    slide_module(prs, "M2", "Gaze & DeepGaze",
                 "Where the user looks. The camera is on the glasses, so head and eye movements still steer what gets stimulated.",
                 ["Heatmaps vs scanpaths; inhibition of return",
                  "DeepGaze III pipeline, plus a synthetic toy model you can play with",
                  "Scanpath sampler with stat histograms",
                  "Why gaze still applies when the eyes still move but can't see",
                  "Owners: Lefteris & Jorge"])
    slide_bullets_image(
        prs, "M2 — Gaze & DeepGaze",
        items=["Heatmaps vs scanpaths; inhibition of return",
               "DeepGaze III pipeline, plus a synthetic toy model you can play with",
               "Scanpath sampler with stat histograms",
               "Why gaze still applies when the eyes move but can't see",
               "Owners: Lefteris & Jorge"],
        image_path=SHOTS_DIR / "M2.png",
        subhead="Where the user looks. The camera is on the glasses, so head and eye movements still steer what gets stimulated.",
        credit="screenshot from modules/M2-deepgaze-and-gaze.html",
        image_h=Inches(2.8),
        qr_path=QRS_DIR / "M2.png",
        qr_label="scan ▶ M2 live",
    )

    # --- M3 ---
    slide_bullets_image(
        prs, "M3 — Neuromodulation & stimulation",
        items=["Biphasic pulse explorer: amp, width, frequency, train shape",
               "Utah array config table you draft and add to",
               "Conductor view: Utah flashing, channels×time, live safety chips",
               "Configure → connect → stim"],
        image="s55_p03_1ff3d95a.png",
        subhead="From a clean visual feature to a safe pulse train on an electrode.",
        credit="Granley & Beyeler — temporal microstim patterns (constant / ramp / biomimetic)",
        image_h=Inches(2.0),
    )
    slide_bullets_image(
        prs, "M3 — Neuromodulation & stimulation",
        items=["Biphasic pulse explorer: amp, width, frequency, train shape",
               "Utah array config table you draft and add to",
               "Conductor view: Utah flashing, channels×time, live safety chips",
               "Configure → connect → stim"],
        image_path=SHOTS_DIR / "M3.png",
        subhead="From a clean visual feature to a safe pulse train on an electrode.",
        credit="screenshot from modules/M3-neuromod-and-stim.html",
        image_h=Inches(2.8),
        qr_path=QRS_DIR / "M3.png",
        qr_label="scan ▶ M3 live",
    )

    # --- M4 ---
    slide_bullets_image(
        prs, "M4 — Phosphene simulation",
        items=["Single-phosphene explorer to feel the basis function",
               "Population viewer with a layout selector",
               "Full image → phosphenes demo, with an animate-drift toggle",
               "Temporal dynamics: leaky integrator + adaptation trace"],
        image="s47_p01_602a4e0f.png",
        subhead="The forward model: stim in, phosphenes out. dynaphos, in the browser.",
        credit="van der Grinten et al. 2024, eLife — dynaphos forward model (CC BY)",
        image_h=Inches(2.4),
    )
    slide_bullets_image(
        prs, "M4 — Phosphene simulation",
        items=["Single-phosphene explorer to feel the basis function",
               "Population viewer with a layout selector",
               "Full image → phosphenes demo, with an animate-drift toggle",
               "Temporal dynamics: leaky integrator + adaptation trace"],
        image_path=SHOTS_DIR / "M4.png",
        subhead="The forward model: stim in, phosphenes out. dynaphos, in the browser.",
        credit="screenshot from modules/M4-phosphene-simulation.html",
        image_h=Inches(2.8),
        qr_path=QRS_DIR / "M4.png",
        qr_label="scan ▶ M4 live",
    )

    # --- M5 ---
    slide_bullets_image(
        prs, "M5 — Decoding & closed loop",
        items=["Closed-loop pipeline with a live mean-pixel brightness readout",
               "Classical PID on dynaphos's leaky integrator, with Kp/Ki/Kd sliders",
               "TF.js MLP trained on synthetic phosphene canvases",
               "Hand-tuned vs end-to-end preprocessor, side by side",
               "Live 2×2 quad: the full loop with an open/closed toggle"],
        image="s49_p01_b5b7a7cf.png",
        subhead="Encoder → simulator → phosphenes → decoder. Reconstruction loss closes the loop end to end.",
        credit="de Ruyter van Steveninck et al. — end-to-end differentiable prosthesis pipeline",
        image_h=Inches(2.3),
    )
    slide_bullets_image(
        prs, "M5 — Decoding & closed loop",
        items=["Closed-loop pipeline with a live mean-pixel brightness readout",
               "Classical PID on dynaphos's leaky integrator, with Kp/Ki/Kd sliders",
               "TF.js MLP trained on synthetic phosphene canvases",
               "Hand-tuned vs end-to-end preprocessor, side by side",
               "Live 2×2 quad: the full loop with an open/closed toggle"],
        image_path=SHOTS_DIR / "M5.png",
        subhead="Encoder → simulator → phosphenes → decoder. Reconstruction loss closes the loop end to end.",
        credit="screenshot from modules/M5-decoding-and-closed-loop.html",
        image_h=Inches(2.6),
        qr_path=QRS_DIR / "M5.png",
        qr_label="scan ▶ M5 live",
    )

    # 6 — Where the field is heading + closing
    slide_image_focus(
        prs, "Where this is heading",
        image="s57_p01_783fa3be.png",
        caption="Next-gen vision implants. Every module today is one piece of a much bigger pipeline.",
        credit="Chris Klink — advanced-pipelines schematic, NIN",
        max_h=Inches(4.6),
    )
    # Team credits — two slides of circular portraits + names. Sources are the
    # same images the prelude deck uses; centralised under build/assets/team/.
    slide_team(
        prs, "NTH team",
        members=[
            ("Samantha Wolff",       TEAM_DIR / "samantha_wolff.png"),
            ("Francesc Varkevisser", TEAM_DIR / "francesc_varkevisser.png"),
            ("Patricija Burgar",     TEAM_DIR / "patricija_burgar.png"),
            ("Antonio Lozano",       TEAM_DIR / "antonio_lozano.png"),
            ("Stijn Balk",           TEAM_DIR / "stijn_balk.png"),
        ],
        cols=3,
    )
    slide_team(
        prs, "Organizing team",
        members=[
            ("Lefteris Papadopoulos", TEAM_DIR / "lefteris_papadopoulos.png"),
            ("Jorge Sanmartin",       TEAM_DIR / "jorge_sanmartin.png"),
            ("Radovan Vodila",        TEAM_DIR / "radovan_vodila.jpg"),
            ("Milan ten Bosch",       TEAM_DIR / "milan_ten_bosch.jpg"),
        ],
        cols=4,
    )

    # (Groups, tracks, prizes, GitHub already covered in the prelude.)
    slide_closing(prs)

    # AIMD highlights — appendix slide after the call-to-action. Names and
    # contributions are first drafts; RV to curate during AIMD prep.
    slide_bullets(
        prs, "AIMD highlights",
        subhead="Vision-restoration contributions worth flagging at AIMD — RV to curate.",
        items=[
            "Eduardo Fernández — UMH Elche · CORTIVIS first-in-human intracortical visual prosthesis (NEJM 2021 Gómez et al.; Sci Adv 2025)",
            "Nick Ramsey — UMC Utrecht · long-term cortical BCI for locked-in communication (NEJM 2016 Vansteensel et al.) and motor decoding",
            "Bert Lourens — [contribution + affiliation TBC by RV]",
        ],
    )

    _stamp_slide_numbers(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    repo = Path(__file__).resolve().parent.parent
    out = repo / "presentations" / "bootcamp_talk_AL.pptx"
    saved = build(out)
    print(f"wrote {saved.relative_to(repo)}  ({saved.stat().st_size:,} bytes)")
