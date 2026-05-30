"""Extract images + text transcript from a .pptx source.

Usage:
    python build/_extract_pptx_assets.py <source.pptx> <output_dir>

Writes:
    <output_dir>/images/s{slide}_p{n}_{hash}.{ext}
    <output_dir>/transcript.txt   — slide-by-slide text dump
    <output_dir>/manifest.tsv     — one row per image: slide, name, w_px, h_px, bytes
"""

from __future__ import annotations

import hashlib
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def walk_shapes(shape, sink):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            walk_shapes(child, sink)
        return
    sink.append(shape)


def extract(src: Path, out: Path) -> None:
    out.mkdir(parents=True, exist_ok=True)
    img_dir = out / "images"
    img_dir.mkdir(exist_ok=True)

    prs = Presentation(str(src))
    transcript_lines: list[str] = []
    manifest: list[str] = ["slide\tfile\twidth_emu\theight_emu\tbytes\tsha8"]

    for s_idx, slide in enumerate(prs.slides, start=1):
        transcript_lines.append(f"\n=== Slide {s_idx} ===")
        shapes: list = []
        for sh in slide.shapes:
            walk_shapes(sh, shapes)

        pic_n = 0
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic_n += 1
                blob = sh.image.blob
                ext = sh.image.ext
                h = hashlib.sha1(blob).hexdigest()[:8]
                fname = f"s{s_idx:02d}_p{pic_n:02d}_{h}.{ext}"
                (img_dir / fname).write_bytes(blob)
                manifest.append(
                    f"{s_idx}\t{fname}\t{sh.width}\t{sh.height}\t{len(blob)}\t{h}"
                )
            elif sh.has_text_frame:
                text = sh.text_frame.text.strip()
                if text:
                    transcript_lines.append(text)

        # also pull slide notes if any
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                transcript_lines.append(f"[notes] {notes_text}")

    (out / "transcript.txt").write_text(
        "\n".join(transcript_lines), encoding="utf-8"
    )
    (out / "manifest.tsv").write_text(
        "\n".join(manifest), encoding="utf-8"
    )

    print(f"slides:     {len(prs.slides)}")
    print(f"images:     {sum(1 for _ in img_dir.iterdir())}")
    print(f"transcript: {(out / 'transcript.txt').relative_to(out.parent)}")
    print(f"manifest:   {(out / 'manifest.tsv').relative_to(out.parent)}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print(__doc__)
        sys.exit(1)
    extract(Path(sys.argv[1]), Path(sys.argv[2]))
